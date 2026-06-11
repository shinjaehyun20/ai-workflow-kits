"""Dev helper: build a donor 'design-system.pptx' to exercise the engine.

This is the ONLY place python-pptx is used, and only to mint a test fixture.
The workspace engine itself stays pure standard library.
"""

from __future__ import annotations

import sys

from pptx import Presentation


def build(path: str) -> str:
    prs = Presentation()  # default template: one master, 11 named layouts
    s1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    s1.shapes.title.text = "로그인 화면"
    s1.placeholders[1].text = "디자인시스템 v1"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    s2.shapes.title.text = "로그인 폼"
    s2.placeholders[1].text = "이메일과 비밀번호를 입력하세요"
    prs.save(path)
    return path


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "design-system.pptx"
    print("wrote", build(out))
